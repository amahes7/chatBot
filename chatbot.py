import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from langchain.chains import ConversationalRetrievalChain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.chat_models import ChatOpenAI
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from openai import OpenAIError
from pptx import Presentation

# ------------------ Read API key ------------------
try:
    with open("OPENAI_API_KEY.txt", "r") as f:
        OPENAI_API_KEY = f.read().strip()
except FileNotFoundError:
    st.error("❌ API key file 'OPENAI_API_KEY.txt' not found.")
    OPENAI_API_KEY = None

# ------------------ Streamlit UI ------------------
st.set_page_config(page_title="🤖 Smart Chatbot", page_icon="🤖", layout="wide")

st.markdown("<h1 style='text-align: center;'>🤖 Smart Chatbot</h1>", unsafe_allow_html=True)
st.markdown("<p style='text-align: center;'>Upload documents (PDF, PPTX, DOCX, TXT) for refined answers</p>", unsafe_allow_html=True)
st.markdown("---")

with st.sidebar:
    st.title("📂 Upload Document")
    file = st.file_uploader(
        "Supported formats: PDF, PPTX, DOCX, TXT",
        type=["pdf", "pptx", "docx", "txt"]
    )

# ------------------ Chat history ------------------
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# Add a system message when a new file is uploaded
if file is not None:
    uploaded_file_message = f"📂 Document uploaded: {file.name}"
    if len(st.session_state.chat_history) == 0 or st.session_state.chat_history[-1][0] != uploaded_file_message:
        st.session_state.chat_history.append((uploaded_file_message, "", []))

# ------------------ Text area + submit button ------------------
with st.form(key="qa_form"):
    user_question = st.text_area("💬 Ask a question:", height=100)
    submitted = st.form_submit_button("Submit")

# ------------------ File text extraction ------------------
def extract_text(file, file_type):
    text = ""
    if file_type == "pdf":
        pdf_reader = PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
    elif file_type == "pptx":
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file_type == "docx":
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_type == "txt":
        text = file.read().decode("utf-8")
    return text

# ------------------ LLM setup ------------------
if OPENAI_API_KEY:
    llm = ChatOpenAI(
        openai_api_key=OPENAI_API_KEY,
        temperature=0,
        max_tokens=1000,
        model_name="gpt-3.5-turbo"
    )

    if submitted and user_question.strip():
        if file is not None:
            file_type = file.name.split(".")[-1].lower()
            text = extract_text(file, file_type)

            if text.strip():
                # Split into chunks
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=150,
                    separators=["\n"],
                    length_function=len
                )
                chunks = text_splitter.split_text(text)

                # Embeddings + Vector store
                embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
                vector_store = FAISS.from_texts(chunks, embeddings)

                # Conversational retrieval chain
                qa_chain = ConversationalRetrievalChain.from_llm(
                    llm=llm,
                    retriever=vector_store.as_retriever(search_kwargs={"k": 3}),
                    return_source_documents=True
                )

                try:
                    result = qa_chain(
                        {"question": user_question, "chat_history": st.session_state.chat_history}
                    )
                    answer = result["answer"]
                    st.session_state.chat_history.append((user_question, answer, result["source_documents"]))

                    with st.expander("🤖 View Answer", expanded=True):
                        st.write(answer)

                except OpenAIError as e:
                    st.error(f"❌ Error: {e}")
            else:
                st.warning("⚠️ No readable text found in this file.")

        else:
            # No file uploaded → fallback LLM answer
            try:
                response = llm.predict(user_question)
                st.session_state.chat_history.append((user_question, response, []))

                with st.expander("🤖 View Answer", expanded=True):
                    st.write(response)

            except OpenAIError as e:
                st.error(f"❌ Error: {e}")

# ------------------ Display conversation history ------------------
if st.session_state.chat_history:
    st.markdown("### 📝 Conversation History")
    for i, (q, a, sources) in enumerate(st.session_state.chat_history, 1):
        if a == "":  # System message (file upload)
            st.markdown(f"**{q}**")
        else:
            with st.expander(f"**Q:** {q}"):
                st.markdown(f"**🤖 Answer:** {a}")
                if sources:
                    with st.expander("📌 Sources"):
                        for doc in sources:
                            st.markdown(doc.page_content[:300] + "...")
